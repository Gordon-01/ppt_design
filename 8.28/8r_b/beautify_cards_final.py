#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
第12、13页卡片化最终版
"""
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = '待美化_第12-13页(workbuddy).pptx'
OUT = '待美化_第12-13页(workbuddy).pptx'

# 配色
C_DEEP_RED = RGBColor(0xC0, 0x00, 0x00)
C_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BORDER = RGBColor(0xE0, 0xE0, 0xE0)

# 字号
SZ_SUBTITLE = Pt(20)
SZ_LABEL = Pt(20)
SZ_BODY = Pt(16)

# 内容
contents = [
    {
        'subtitle': '过期标语未及时撤除，易造成经营风险',
        'label1': '问题点：',
        'body1': '过期标语仍在市场运用（提神快、十包高端槟榔七包和成天下），未及时撤除。其中偃师市场情况较为严重，包含陈列盒/架以及终端品宣。',
        'label2': '改善建议：',
        'body2': '1.市场全面大清洗，动员市场所有人员全面撤除过期品宣；\n2.对人员密集、地理位置优越（景区、十字路口、步行街等）的终端给予大品宣费用支持；',
        'has_photos': True,
    },
    {
        'subtitle': '骄阳行动颁奖进度缓慢，线上舆论增加',
        'label1': '问题点：',
        'body1': '奖品（电动车、手机）到达市场时间较长，市场颁奖进度缓慢，导致中奖客户在线上发布不实信息，给予市场较大压力。',
        'label2': '改善建议：',
        'body2': '1.确定市场奖品数量并提前发往市场，保证颁奖质量的同时减少复杂流程；\n2.制定相关话术，业务员灵活使用话术安抚中奖客户情绪，减少相关舆论，维护本品形象；',
        'has_photos': False,
    },
]


def set_run_font(run, name, size, bold, color):
    run.font.name = name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def set_paragraph_text(para, text, name, size, bold, color, align=None, line_spacing=1.2, space_after=0):
    para.clear()
    run = para.add_run()
    run.text = text
    set_run_font(run, name, size, bold, color)
    if align is not None:
        para.alignment = align
    para.line_spacing = line_spacing
    para.space_after = Pt(space_after)
    para.space_before = Pt(0)


def add_card(slide, x, y, w, h, label_text, body_text):
    # 卡片底
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(x), int(y), int(w), int(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG
    card.line.color.rgb = C_CARD_BORDER
    card.line.width = Pt(1.0)

    # 左侧红色竖条
    accent_w = Inches(0.10)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(x), int(y), int(accent_w), int(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_DEEP_RED
    bar.line.fill.background()

    # 文字区
    text_left = int(x + accent_w + Inches(0.14))
    text_top = int(y + Inches(0.13))
    text_w = int(w - accent_w - Inches(0.28))
    text_h = int(h - Inches(0.26))
    tb = slide.shapes.add_textbox(text_left, text_top, text_w, text_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    # 标签
    set_paragraph_text(
        tf.paragraphs[0], label_text,
        '微软雅黑', SZ_LABEL, True, C_DEEP_RED,
        align=None, line_spacing=1.2, space_after=6
    )

    # 正文
    p2 = tf.add_paragraph()
    set_paragraph_text(
        p2, body_text,
        '微软雅黑', SZ_BODY, False, C_DARK_GRAY,
        align=None, line_spacing=1.35, space_after=0
    )

    return card, bar, tb


prs = Presentation(SRC)
assert len(prs.slides) == 2

for slide_idx, slide in enumerate(prs.slides):
    info = contents[slide_idx]

    # 1) 格式化副标题条：整段替换为单一 run，避免重复
    for sh in slide.shapes:
        if sh.name == '矩形 5':
            tf = sh.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            for para in tf.paragraphs:
                set_paragraph_text(
                    para, info['subtitle'],
                    '微软雅黑', SZ_SUBTITLE, True, C_WHITE,
                    align=PP_ALIGN.CENTER, line_spacing=1.2, space_after=0
                )
            break

    # 2) 收集内容照片（排除模板 logo / 装饰）
    photos = []
    for sh in slide.shapes:
        if sh.shape_type == 13 and sh.name not in ('图片 7', '图片 101'):
            photos.append(sh)

    # 3) 删除旧文本框
    for sh in list(slide.shapes):
        if sh.name in ('文本框 6', '文本框 8'):
            sp = sh._element
            sp.getparent().remove(sp)

    # 4) 计算卡片布局
    slide_w = prs.slide_width
    margin_x = Inches(0.55)
    content_w = slide_w - margin_x * 2
    card_x = margin_x
    card_w = content_w

    content_top = Inches(2.42)
    content_bottom = Inches(6.45)
    content_h = content_bottom - content_top

    if info['has_photos'] and photos:
        photo_h = Inches(1.05)
        photo_gap = Inches(0.22)
        gap_card_photo = Inches(0.18)
        gap_photo_card = Inches(0.18)
        card_h = (content_h - photo_h - gap_card_photo - gap_photo_card) / 2
        card1_y = content_top
        photo_y = card1_y + card_h + gap_card_photo
        card2_y = photo_y + photo_h + gap_photo_card
    else:
        gap = Inches(0.28)
        card_h = (content_h - gap) / 2
        card1_y = content_top
        card2_y = card1_y + card_h + gap
        photo_y = None

    # 5) 添加卡片
    add_card(slide, card_x, card1_y, card_w, card_h, info['label1'], info['body1'])
    add_card(slide, card_x, card2_y, card_w, card_h, info['label2'], info['body2'])

    # 6) 重新排列照片
    if photo_y is not None and photos:
        # 按原文件名排序，保持视觉顺序
        photos_sorted = sorted(photos, key=lambda s: s.name)
        photo_items = []
        for sh in photos_sorted:
            aspect = sh.width / sh.height if sh.height else 1.0
            photo_items.append((sh, aspect))
            if photo_items:
                total_w = sum(photo_h * aspect for _, aspect in photo_items)
                total_gap = photo_gap * (len(photo_items) - 1)
                group_w = total_w + total_gap
                # 照片组整体居中
                start_x = margin_x + (content_w - group_w) / 2
                cur_x = start_x
                for sh, aspect in photo_items:
                    new_w = photo_h * aspect
                    sh.top = int(photo_y)
                    sh.left = int(cur_x)
                    sh.width = int(new_w)
                    sh.height = int(photo_h)
                    cur_x += int(new_w + photo_gap)

prs.save(OUT)
print(f"Saved to {OUT}")
