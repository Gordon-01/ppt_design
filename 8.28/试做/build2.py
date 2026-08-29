from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement
import os

# ================= 设计系统 =================
A = "assets2"
OUT = "我的暑假生活.pptx"
FONT = "Microsoft YaHei"

SW, SH = Inches(13.333), Inches(7.5)
M = Inches(0.75)                 # 左右边距
CL, CR = Inches(0.75), Inches(12.583)   # 内容区左右边界
CONTENT_W = CR - CL              # 11.833"
HEADER_Y, HEADER_H = Inches(0.52), Inches(0.78)
DIVIDER_Y = Inches(1.55)
BODY_T, BODY_B = Inches(1.85), Inches(6.85)

BG      = RGBColor(0xFF, 0xFD, 0xF9)
TEXT    = RGBColor(0x2B, 0x34, 0x40)
MUTED   = RGBColor(0x8A, 0x93, 0xA0)
LINE    = RGBColor(0xE8, 0xEA, 0xEE)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x4F, 0xA8, 0xE8)
CORAL   = RGBColor(0xFF, 0x8A, 0x7A)
LAV     = RGBColor(0xB3, 0x9D, 0xDB)
SUN     = RGBColor(0xFF, 0xB8, 0x4D)
MINT    = RGBColor(0x5F, 0xCF, 0xA8)
INK     = RGBColor(0x1F, 0x29, 0x37)


def sf(run, size, bold=False, color=TEXT, name=FONT):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = name


def txt(slide, l, t, w, h, s, size, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=0, line_spacing=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = s.split("\n") if isinstance(s, str) else s
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.space_after = Pt(spacing)
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run(); r.text = line; sf(r, size, bold, color)
    return box


def rect(slide, l, t, w, h, fill=WHITE, line=None, lw=Pt(1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1, shadow=False):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    try:
        if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            s.adjustments[0] = radius
    except Exception:
        pass
    if shadow:
        add_shadow(s)
    return s


def add_shadow(shape, blur=152400, dist=38100, alpha=11000):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    eff = parse_xml(
        '<a:effectLst %s><a:outerShdw blurRad="%d" dist="%d" dir="5400000" rotWithShape="0">'
        '<a:srgbClr val="1F2937"><a:alpha val="%d"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>' % (nsdecls('a'), blur, dist, alpha))
    spPr.append(eff)


def set_alpha(shape, val):
    """给纯色填充加透明度，val 为千分比百分比，如 35000 = 35%"""
    spPr = shape._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None: return
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None: return
    for a in srgb.findall(qn('a:alpha')):
        srgb.remove(a)
    al = OxmlElement('a:alpha'); al.set('val', str(val))
    srgb.append(al)


def pic(slide, path, l, t, w, h, radius=0.1, shadow=False):
    """圆角图片：用 rounded rect + 图片填充实现"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try: s.adjustments[0] = radius
    except Exception: pass
    s.line.fill.background()
    img_part, rId = slide.part.get_or_add_image_part(path)
    spPr = s._element.spPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    blip = parse_xml(
        '<a:blipFill %s %s rotWithShape="1">'
        '<a:blip r:embed="%s"/><a:stretch><a:fillRect/></a:stretch>'
        '</a:blipFill>' % (nsdecls('a'), nsdecls('r'), rId))
    geom = spPr.find(qn('a:prstGeom'))
    if geom is not None: geom.addnext(blip)
    else: spPr.insert(0, blip)
    if shadow: add_shadow(s)
    return s


def header(slide, num, title, accent):
    chip = rect(slide, CL, HEADER_Y, Inches(0.78), HEADER_H, fill=accent, radius=0.22)
    tf = chip.text_frame; tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; sf(r, 24, True, WHITE)
    txt(slide, CL + Inches(0.97), HEADER_Y - Inches(0.03), Inches(7), HEADER_H,
        title, 34, True, TEXT, anchor=MSO_ANCHOR.MIDDLE)
    rule = rect(slide, CL, DIVIDER_Y, CONTENT_W, Pt(1), fill=LINE,
                shape=MSO_SHAPE.RECTANGLE, radius=0)
    return chip


def footer(slide, page, total=6):
    txt(slide, CR - Inches(1.2), Inches(6.95), Inches(1.2), Inches(0.32),
        f"{page} / {total}", 11, False, MUTED, align=PP_ALIGN.RIGHT)
    txt(slide, CL, Inches(6.95), Inches(5), Inches(0.32),
        "我的暑假生活", 11, False, MUTED)


def bullet_cards(slide, items, x, y, w, accent, card_h=Inches(1.02),
                 gap=Inches(0.2), size=17):
    for i, line in enumerate(items):
        top = y + i * (card_h + gap)
        card = rect(slide, x, top, w, card_h, fill=WHITE, radius=0.16, shadow=True)
        bar = rect(slide, x + Inches(0.2), top + Inches(0.24), Inches(0.09),
                   card_h - Inches(0.48), fill=accent, radius=0.5)
        txt(slide, x + Inches(0.5), top, w - Inches(0.7), card_h,
            line, size, False, TEXT, anchor=MSO_ANCHOR.MIDDLE)
    return y + len(items) * (card_h + gap) - gap


# ================= 生成 =================
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
blank = prs.slide_layouts[6]

# ---------- 1. 封面 ----------
s = prs.slides.add_slide(blank)
LEFT_W = Inches(5.4)
# 左侧暖色面板
rect(s, 0, 0, LEFT_W, SH, fill=RGBColor(0xFF, 0xF6, 0xE8), shape=MSO_SHAPE.RECTANGLE, radius=0)
# 右侧满幅图片
pic(s, f"{A}/cover_split.jpg", LEFT_W, 0, SW - LEFT_W, SH, radius=0.0)
# 装饰元素
rect(s, Inches(0.95), Inches(2.35), Inches(1.05), Inches(0.13), fill=SUN, radius=0.5)
txt(s, Inches(0.95), Inches(2.62), Inches(4.2), Inches(1.25),
    "我的暑假生活", 62, True, TEXT)
txt(s, Inches(0.98), Inches(4.0), Inches(4.2), Inches(0.55),
    "姓名：XXX        班级：X年X班", 24, False, TEXT)
chip = rect(s, Inches(0.95), Inches(4.85), Inches(2.65), Inches(0.58), fill=SUN, radius=0.3)
tf = chip.text_frame; tf.margin_left = 0; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "2026 · 暑期生活分享"; sf(r, 15, True, INK)

# ---------- 2. 自律学习 ----------
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
header(s, "01", "自律学习", BLUE)
bullet_cards(s, [
    "制定暑假学习计划，劳逸结合",
    "每天认真完成假期作业",
    "坚持阅读课外书，完成阅读集章卡任务",
    "每天练字，让字迹更工整",
], CL, Inches(1.95), Inches(6.55), BLUE, size=17)
pic(s, f"{A}/study.jpg", Inches(7.98), BODY_T, Inches(4.6), Inches(5.0),
    radius=0.09, shadow=True)
rect(s, Inches(7.78), Inches(6.62), Inches(1.7), Inches(0.14), fill=SUN, radius=0.5)
footer(s, 2)

# ---------- 3. 锻炼与新本领 ----------
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
header(s, "02", "锻炼与新本领", CORAL)
cards = [
    (f"{A}/swim.jpg", CORAL, "游泳小健将",
     ["坚持每天体育锻炼", "学会自由泳和跳水", "挑战自己，变得更勇敢"]),
    (f"{A}/housework.jpg", MINT, "家务小帮手",
     ["主动帮爸爸妈妈做家务", "学会整理、打扫和简单烹饪", "懂得劳动的意义，更有担当"]),
]
cw = Inches(5.61); ch = Inches(4.95); cy = Inches(1.85); gapx = Inches(0.613)
for i, (img, accent, title, bullets) in enumerate(cards):
    cx = CL + i * (cw + gapx)
    rect(s, cx, cy, cw, ch, fill=WHITE, radius=0.075, shadow=True)
    iw = cw - Inches(0.32); ih = Inches(2.97)
    pic(s, img, cx + Inches(0.16), cy + Inches(0.16), iw, ih, radius=0.06)
    tag = rect(s, cx + Inches(0.34), cy + Inches(0.16) + ih - Inches(0.62),
               Inches(2.5), Inches(0.46), fill=accent, radius=0.22)
    tf = tag.text_frame; tf.margin_left = 0; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; sf(r, 17, True, WHITE)
    for j, b in enumerate(bullets):
        by = cy + Inches(0.16) + ih + Inches(0.34) + j * Inches(0.45)
        dot = rect(s, cx + Inches(0.34), by + Inches(0.11), Inches(0.16),
                   Inches(0.16), fill=accent, radius=0.5)
        txt(s, cx + Inches(0.62), by, cw - Inches(0.9), Inches(0.42),
            b, 15, False, TEXT, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# ---------- 4. 社会实践 ----------
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
header(s, "03", "社会实践", LAV)
pic(s, f"{A}/bookstore.jpg", CL, BODY_T, Inches(5.6), Inches(5.0),
    radius=0.09, shadow=True)
rect(s, CL + Inches(0.25), Inches(6.62), Inches(1.7), Inches(0.14), fill=LAV, radius=0.5)
bullet_cards(s, [
    "在新华书店体验一周小店员",
    "整理书架、接待读者，体会工作的辛苦",
    "参加两期阅读公益讲堂",
    "学到许多书本以外的知识",
], Inches(6.9), Inches(1.95), Inches(5.68), LAV, size=16)
footer(s, 4)

# ---------- 5. 快乐旅行 ----------
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
header(s, "04", "快乐旅行", SUN)
txt(s, CL, Inches(1.78), Inches(10), Inches(0.4),
    "和妈妈一起游览河西走廊与西安，欣赏祖国美景与悠久历史文化。", 19, False, TEXT)
spots = [("danxia.jpg", "七彩丹霞"), ("mogao.jpg", "莫高窟"),
         ("greatwall.jpg", "长城"), ("terracotta.jpg", "兵马俑")]
cw2 = Inches(2.70); gapx2 = Inches(0.343); cy2 = Inches(2.35)
ih2 = Inches(2.025)
for i, (f, label) in enumerate(spots):
    cx = CL + i * (cw2 + gapx2)
    pic(s, f"{A}/{f}", cx, cy2, cw2, ih2, radius=0.07, shadow=True)
    txt(s, cx, cy2 + ih2 + Inches(0.16), cw2, Inches(0.35),
        label, 17, True, TEXT, align=PP_ALIGN.CENTER)
band = rect(s, CL, Inches(5.35), CONTENT_W, Inches(0.92),
            fill=RGBColor(0xFF, 0xF6, 0xE5), radius=0.12)
txt(s, CL + Inches(0.4), Inches(5.35), CONTENT_W - Inches(0.8), Inches(0.92),
    "读万卷书，行万里路——这次旅行让我看到祖国的壮美，也感受到历史的厚重。",
    18, False, TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# ---------- 6. 结尾 ----------
s = prs.slides.add_slide(blank)
pic(s, f"{A}/ending.jpg", 0, 0, SW, SH, radius=0.0)
scrim = rect(s, 0, 0, SW, SH, fill=WHITE, shape=MSO_SHAPE.RECTANGLE, radius=0)
set_alpha(scrim, 40000)
cardw, cardh = Inches(8.6), Inches(5.45)
cardx, cardy = (SW - cardw) // 2, Inches(1.0)
rect(s, cardx, cardy, cardw, cardh, fill=WHITE, radius=0.055, shadow=True)
txt(s, cardx, cardy + Inches(0.4), cardw, Inches(0.85),
    "我的收获", 42, True, TEXT, align=PP_ALIGN.CENTER)
rect(s, (SW - Inches(1.2)) // 2, cardy + Inches(1.32), Inches(1.2), Inches(0.1),
     fill=SUN, radius=0.5)
rows = [(MINT, "知识 · 阅读与旅行，让我看到更大的世界"),
        (CORAL, "勇气 · 游泳和跳水，让我敢于挑战自己"),
        (BLUE, "成长 · 劳动与实践，让我学会担当")]
for i, (c, line) in enumerate(rows):
    ry = cardy + Inches(1.72) + i * Inches(0.62)
    rect(s, cardx + Inches(1.35), ry + Inches(0.1), Inches(0.18), Inches(0.18),
         fill=c, radius=0.5)
    txt(s, cardx + Inches(1.72), ry, cardw - Inches(3.0), Inches(0.38),
        line, 21, False, TEXT, anchor=MSO_ANCHOR.MIDDLE)
rect(s, cardx + Inches(1.35), cardy + Inches(3.62), cardw - Inches(2.7), Pt(1),
     fill=LINE, shape=MSO_SHAPE.RECTANGLE, radius=0)
txt(s, cardx, cardy + Inches(3.82), cardw, Inches(0.55),
    "新学期，我会继续努力，做更好的自己！", 24, True, TEXT, align=PP_ALIGN.CENTER)
txt(s, cardx, cardy + Inches(4.4), cardw, Inches(0.9),
    "谢谢大家！", 52, True, CORAL, align=PP_ALIGN.CENTER)
footer(s, 6)

prs.save(OUT)
print("Saved", OUT, os.path.getsize(OUT), "bytes;", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
